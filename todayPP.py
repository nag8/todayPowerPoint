# coding: UTF-8

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from bs4 import BeautifulSoup
import configparser
import requests
import feedparser
import datetime
import csv
import subprocess


# メイン処理
def main():

    print('処理開始！')
    print('自動でファイルと、メモ帳が開きます…')

    # 設定ファイル取得
    iniFile = getIniFile()

    # PowerPointファイルを生成
    createPP(iniFile)

    print('処理終了！')

# 設定ファイル取得
def getIniFile():
    iniFile = configparser.ConfigParser()
    iniFile.read('./config.ini', 'UTF-8')
    return iniFile


# パワーポイントを生成
def createPP(iniFile):
    prs = Presentation(iniFile.get('settings', 'IN'))

    # htmlから情報を取得
    inputTable = getInputTable(iniFile)

    # 表を変更
    editPPTable(iniFile, prs.slides[0].shapes[0].table, prs.slides[1].shapes[1].table, inputTable)

    # ファイルを保存
    prs.save(iniFile.get('settings', 'OUT') + createFileName() + '.pptx')
    subprocess.call("start " + iniFile.get('settings', 'OUT') + createFileName() + '.pptx',shell=True)


# データ元の情報を取得
def getInputTable(iniFile):

    # 保存したhtmlを取得
    with open(iniFile.get('settings', 'HTML'), encoding="shift_JIS", errors='ignore') as f:
        html = f.read()

    #要素を抽出
    soup = BeautifulSoup(html, 'lxml')

    # テーブルを指定
    return soup.findAll("table")[0]

# セルのフォントサイズを変更して、中央揃えにする
def changeLayout(cell, size):
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)

        # 中央揃えにもする
        paragraph.alignment = PP_ALIGN.CENTER


# パワーポイントのテーブルを修正
def editPPTable(iniFile, table1, table2, inputTable):
    # 要素を取得
    tdList = inputTable.findAll("td", attrs = {"class": "p11pa2"})

    # 設定に必要なcsvを取得
    directory = getDirectory(iniFile.get('settings', 'CSV'))



    for td in tdList:

        # 行番号がある場合
        if directory[td.text[:3]][1] != '':

            # 行番号をcsvから取得
            rowNum = int(directory[td.text[:3]][1])

            # 部屋の名前をログ出力
            print('---------------------　-------------')
            print('部屋→→→→→→→→→→    ' + directory[td.text[:3]][0])

            # 要素を取得
            contents = td.parent.findAll("td", attrs = {"class": "p11"})

            for i,content in enumerate(contents):

                # content.replace("®","")
                print(1)

                # テーブル番号、列番号を設定
                if directory[td.text[:3]][2] == '1':
                    changeTable = table1
                    columnNum   = 2 + i
                    pageId = 1

                else:
                    changeTable = table2
                    columnNum   = 1 + i
                    pageId = 2

                # テキストを設定
                changeCell = changeTable.cell(rowNum, columnNum)
                changeCell.text = getStr(content, pageId)

                # エラーになるのでこれだけ除外
                if "®" not in changeCell.text:
                    print(changeCell.text)

                # レイアウトを修正
                changeLayout(changeCell, 10)

    # いつかやる
    # table2.cell(1, 3).merge(table2.cell(2, 3))

# 文字を取得
def getStr(content, pageId):


    strList = content.get_text(';').split(';')

    # 未入金の場合があるので、それを削除
    if strList[0] == '未':
        strList.pop(0)

    nameStr = removeFirstName(strList[0])

    if len(strList) >= 2:
        if pageId == 1:
            return strList[1] + '\n（' + strList[0] + '　様）'
        else:
            return strList[0] + '　様'

    elif len(strList) == 1:
        return strList[0]

    else:
        return 'error'


# CSVを取得
def getDirectory(csvPath):

    directory = {}

    with open(csvPath, 'r') as f:
        reader = csv.reader(f)

        for row in reader:
            directory[row[0]] = [row[1],row[2],row[3]]

    return directory



def removeFirstName(fullName):

    # return fullName.translate(fullNameTable)
    return fullName

# ファイルネームを生成
def createFileName():

    # 曜日
    yobi = ["月","火","水","木","金","土","日"]

    # 明日を取得
    tomorrow = datetime.datetime.now() + datetime.timedelta(days = 1)
    # 整形して返却
    return '{}月{}日（{}）'.format(tomorrow.month, tomorrow.day, yobi[tomorrow.weekday()])



if __name__ == '__main__':
    main()
